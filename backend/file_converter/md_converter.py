import pypandoc
# import docx2md
# from pptx2md import convert, ConversionConfig
import pdfplumber
import tempfile
import pandas as pd
import os
from docx2pdf import convert as cv
from pptx import Presentation


class Txt:
    def __init__(self,file):
        self.file=file
    
    async def convert_to_md (self):
        with open(self.file, "r", encoding="utf-8") as f:
            md_text = f.read()
        converted = pypandoc.convert_text(md_text, to="md", format="markdown",extra_args=["--standalone"])
        return converted
    
class Docx:
    def __init__(self,file):
        self.file = file
    
    async def convert_to_md(self):

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name
        
        try:
            cv(self.file, tmp_path)
            
            md_text = ""
            with pdfplumber.open(tmp_path) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text()
                    if page_text:
                        md_text += f"\n--- Page {i} ---\n{page_text}\n"
            
            return md_text
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

class Ppx:
    def __init__(self,file):
        self.file=file
    
    async def convert_to_md (self):

        prs = Presentation(self.file)        
        md_text = ""
        for slide_num, slide in enumerate(prs.slides, start=1):
            md_text += f"--- Page {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    md_text += f"{shape.text.strip()}\n"
        
        return md_text

class Csv:
    def __init__(self, file):
        self.file = file

    async def convert_to_md(self):
        df = pd.read_csv(self.file)
        md_text = df.to_markdown(tablefmt="grid")
        return md_text
    
class Xls:
    def __init__(self, file):
        self.file = file

    async def convert_to_md(self):
       
        df = pd.read_excel(self.file)
        md_text = df.to_markdown(tablefmt="grid")
        return md_text
    
class Pdf:
    def __init__(self, file):
        self.file = file

    async def convert_to_md(self):
        md_text = ""
        with pdfplumber.open(self.file) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text()
                if page_text:
                    md_text += f"\n--- Page {i} ---\n{page_text}\n"
        return md_text

def get_converter(input_file ,file_name):
    
    ext = file_name.split(".")[-1].lower()

    if ext == "txt":
        return Txt(input_file)
    elif ext == "docx":
        return Docx(input_file)
    elif ext == "pptx":
        return Ppx(input_file)
    elif ext == "csv":
        return Csv(input_file)
    elif ext == "xls" or ext== 'xlsx':
        return Xls(input_file)
    elif ext == "pdf":
        return Pdf(input_file)
    else:
        raise ValueError(f"Unsupported file type: {ext}")